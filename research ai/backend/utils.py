import os
import pandas as pd
from docx import Document
from pptx import Presentation
import PyPDF2  # or pdfplumber if you are currently using that

def extract_text(file_path: str, ext: str):
    """
    Extracts text from various file formats.
    Returns a list of tuples: [(text_chunk, metadata_string), ...]
    """
    raw_data =[]
    
    try:
        if ext == '.pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for i, page in enumerate(reader.pages):
                    text = page.extract_text()
                    if text and text.strip():
                        raw_data.append((text, f"Page {i+1}"))
                        
        elif ext == '.txt':
            # Use errors='ignore' so strange characters don't crash the upload
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                raw_data.append((f.read(), "Text Document"))
                
        elif ext == '.docx':
            doc = Document(file_path)
            full_text =[para.text for para in doc.paragraphs if para.text.strip()]
            raw_data.append(('\n'.join(full_text), "Word Document"))
            
        elif ext == '.pptx':
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_text =[]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    raw_data.append(('\n'.join(slide_text), f"Slide {i+1}"))
                    
        elif ext in ['.xlsx', '.xls']:
            # Read Excel sheets
            df = pd.read_excel(file_path)
            # Convert dataframe to a readable string format for the AI
            raw_data.append((df.to_string(), "Excel Data"))
            
        elif ext == '.csv':
            # Read CSV
            df = pd.read_csv(file_path)
            raw_data.append((df.to_string(), "CSV Data"))
            
        else:
            print(f"Unsupported file type ignored: {ext}")

    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        
    return raw_data

# Keep your existing split_text function down here...
def split_text(text, metadata, filename):
    # Your existing splitting logic goes here...
    pass