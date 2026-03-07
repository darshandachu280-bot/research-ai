import os

# Define the project name
base_dir = "ResearchAI_Final_Project"

# Define the folder structure
structure = [
    "backend",
    "backend/uploads",
    "backend/vectors",
    "frontend",
    "frontend/css"
]

# Define all file contents
files = {
    "requirements.txt": """fastapi
uvicorn
sqlalchemy
python-multipart
python-jose[cryptography]
passlib[bcrypt]
pypdf
sentence-transformers
faiss-cpu
numpy
transformers
torch
python-docx
python-pptx
openpyxl
groq
python-dotenv
fpdf2
python-pptx""",

    ".env": "GROQ_API_KEY=your_key_here",

    "backend/database.py": """from sqlalchemy import create_engine, Column, Integer, String, DateTime, ForeignKey, Text
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker
import datetime
SQLALCHEMY_DATABASE_URL = "sqlite:///./research_assistant.db"
engine = create_engine(SQLALCHEMY_DATABASE_URL, connect_args={"check_same_thread": False, "timeout": 30})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

class User(Base):
    __tablename__ = "users"
    id = Column(Integer, primary_key=True, index=True)
    name = Column(String); email = Column(String, unique=True); hashed_password = Column(String)

class Document(Base):
    __tablename__ = "documents"
    id = Column(Integer, primary_key=True); filename = Column(String); user_id = Column(Integer, ForeignKey("users.id"))
    status = Column(String, default="processing"); uploaded_at = Column(DateTime, default=datetime.datetime.utcnow)

class ChatMessage(Base):
    __tablename__ = "chat_messages"
    id = Column(Integer, primary_key=True); user_id = Column(Integer, ForeignKey("users.id"))
    question = Column(Text); answer = Column(Text); sources = Column(Text); created_at = Column(DateTime, default=datetime.datetime.utcnow)

Base.metadata.create_all(bind=engine)
def get_db():
    db = SessionLocal()
    try: yield db
    finally: db.close()""",

    "backend/auth.py": """from datetime import datetime, timedelta
from jose import JWTError, jwt
from passlib.context import CryptContext
from fastapi import Depends, HTTPException, status
from fastapi.security import OAuth2PasswordBearer
from sqlalchemy.orm import Session
import database

SECRET_KEY = "RESEARCH_AI_EMERALD_SECRET_KEY"
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 1440 
pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="login")

def get_password_hash(password): return pwd_context.hash(password)
def verify_password(plain, hashed): return pwd_context.verify(plain, hashed)
def create_access_token(data: dict):
    to_encode = data.copy()
    to_encode.update({"exp": datetime.utcnow() + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

def get_current_user(token: str = Depends(oauth2_scheme), db: Session = Depends(database.get_db)):
    exc = HTTPException(status_code=401, detail="Unauthorized")
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        email: str = payload.get("sub")
        if email is None: raise exc
    except JWTError: raise exc
    user = db.query(database.User).filter(database.User.email == email).first()
    if user is None: raise exc
    return user""",

    "backend/utils.py": """import PyPDF2, os, openpyxl
from docx import Document as DocxDoc
from pptx import Presentation

def extract_text(file_path, extension):
    chunks = []
    try:
        if extension == '.pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for i, p in enumerate(reader.pages):
                    t = p.extract_text()
                    if t: chunks.append((t.strip(), i + 1))
        elif extension == '.docx':
            doc = DocxDoc(file_path)
            chunks.append((". ".join([p.text for p in doc.paragraphs]), 1))
        elif extension == '.pptx':
            prs = Presentation(file_path)
            for i, s in enumerate(prs.slides):
                t = " ".join([sh.text for sh in s.shapes if hasattr(sh, "text")])
                chunks.append((t, i + 1))
        elif extension == '.xlsx':
            wb = openpyxl.load_workbook(file_path, data_only=True)
            for s in wb.worksheets:
                r = [" ".join([str(c) for c in row if c]) for row in s.iter_rows(values_only=True)]
                chunks.append((". ".join(r), 1))
        elif extension == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                chunks.append((f.read(), 1))
    except Exception as e: print(f"Parser Error: {e}")
    return chunks

def split_text(text, page_num, filename):
    w = text.split()
    return [{"text": " ".join(w[i:i+500]), "page": page_num, "filename": filename} for i in range(0, len(w), 450)]""",

    "backend/rag_engine.py": """import os, pickle, faiss, numpy as np
from sentence_transformers import SentenceTransformer
from groq import Groq
def get_client(): return Groq(api_key=os.environ.get("GROQ_API_KEY"))
print("AI: Initializing Embedding Engine..."); embed_model = SentenceTransformer('all-MiniLM-L6-v2')
def build_vector_store(uid, chunks):
    if not chunks: return False
    try:
        texts = [c['text'] for c in chunks]; embs = embed_model.encode(texts)
        idx = faiss.IndexFlatL2(embs.shape[1]); idx.add(np.array(embs).astype('float32'))
        os.makedirs("vectors", exist_ok=True)
        faiss.write_index(idx, f"vectors/u_{uid}.index")
        with open(f"vectors/u_{uid}.pkl", "wb") as f: pickle.dump(chunks, f)
        return True
    except: return False
def get_answer(uid, q, focus_file=None):
    if not os.path.exists(f"vectors/u_{uid}.index"): return "No docs.", []
    idx = faiss.read_index(f"vectors/u_{uid}.index")
    with open(f"vectors/u_{uid}.pkl", "rb") as f: meta = pickle.load(f)
    vec = embed_model.encode([q]).astype('float32'); D, I = idx.search(vec, k=15)
    ctx = []
    for i in I[0]:
        if i != -1 and i < len(meta):
            if focus_file and meta[i]['filename'] != focus_file: continue
            ctx.append(meta[i])
    ctx_text = "\\n\\n".join([f"Doc: {c['filename']} (p{c['page']})\\n{c['text']}" for c in ctx[:6]])
    try:
        res = get_client().chat.completions.create(
            messages=[{"role":"system","content":"Answer strictly using context. Use Markdown. If not found, say 'Not found in document.' At the end add 2 follow-ups starting with 'SUGGESTED:'"},{"role":"user","content":f"CONTEXT:\\n{ctx_text}\\n\\nQ: {q}"}],
            model="llama-3.1-8b-instant", temperature=0.1)
        return res.choices[0].message.content, ctx[:6]
    except Exception as e: return f"AI Error: {e}", []
def generate_document_summary(uid):
    if not os.path.exists(f"vectors/u_{uid}.pkl"): return "No documents."
    with open(f"vectors/u_{uid}.pkl", "rb") as f: meta = pickle.load(f)
    text = " ".join([c['text'] for c in meta[:10]])
    try:
        res = get_client().chat.completions.create(
            messages=[{"role":"system","content":"Summarize the research highlights in bullet points."},{"role":"user","content":text[:7000]}],
            model="llama-3.1-8b-instant", temperature=0.3)
        return res.choices[0].message.content
    except Exception as e: return f"Error: {e}" """,

    "backend/main.py": """import database, auth, utils, rag_engine, json, os, shutil
from dotenv import load_dotenv
load_dotenv(); os.environ["GROQ_API_KEY"] = os.getenv("GROQ_API_KEY", "")
from typing import List, Optional
from fastapi import FastAPI, Depends, UploadFile, File, Form, BackgroundTasks, HTTPException
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from sqlalchemy.orm import Session
from fpdf import FPDF
app = FastAPI()
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

def run_proc(uid, db_factory):
    db = db_factory()
    try:
        docs = db.query(database.Document).filter(database.Document.user_id==uid, database.Document.status=="processing").all()
        for d in docs: d.status = "completed"
        db.commit()
        all_docs = db.query(database.Document).filter(database.Document.user_id==uid, database.Document.status=="completed").all()
        chunks = []
        for d in all_docs:
            raw = utils.extract_text(f"uploads/{uid}_{d.filename}", os.path.splitext(d.filename)[1].lower())
            for t, p in raw: chunks.extend(utils.split_text(t, p, d.filename))
        if chunks: rag_engine.build_vector_store(uid, chunks)
    finally: db.close()

@app.post("/register")
def reg(name:str=Form(...), email:str=Form(...), password:str=Form(...), db:Session=Depends(database.get_db)):
    if db.query(database.User).filter(database.User.email == email).first(): raise HTTPException(400)
    db.add(database.User(name=name, email=email, hashed_password=auth.get_password_hash(password)))
    db.commit(); return {"msg": "ok"}
@app.post("/login")
def login(email:str=Form(...), password:str=Form(...), db:Session=Depends(database.get_db)):
    u = db.query(database.User).filter(database.User.email==email).first()
    if not u or not auth.verify_password(password, u.hashed_password): raise HTTPException(401)
    return {"access_token": auth.create_access_token({"sub": u.email}), "token_type": "bearer"}
@app.get("/me")
def me(u: database.User = Depends(auth.get_current_user)): return u
@app.post("/upload")
def upload(bt: BackgroundTasks, files: List[UploadFile] = File(...), u=Depends(auth.get_current_user), db: Session = Depends(database.get_db)):
    os.makedirs("uploads", exist_ok=True)
    for f in files:
        path = f"uploads/{u.id}_{f.filename}"
        with open(path, "wb") as b: shutil.copyfileobj(f.file, b)
        db.add(database.Document(filename=f.filename, user_id=u.id, status="processing"))
    db.commit(); bt.add_task(run_proc, u.id, database.SessionLocal); return {"msg": "ok"}
@app.get("/documents")
def docs(u: database.User = Depends(auth.get_current_user), db: Session = Depends(database.get_db)):
    return db.query(database.Document).filter(database.Document.user_id == u.id).all()
@app.delete("/documents/{doc_id}")
def delete_doc(doc_id: int, bt: BackgroundTasks, u: database.User = Depends(auth.get_current_user), db: Session = Depends(database.get_db)):
    doc = db.query(database.Document).filter(database.Document.id == doc_id, database.Document.user_id == u.id).first()
    if doc:
        p = f"uploads/{u.id}_{doc.filename}"
        if os.path.exists(p): os.remove(p)
        db.delete(doc); db.commit(); bt.add_task(run_proc, u.id, database.SessionLocal)
    return {"msg": "ok"}
@app.post("/ask")
def ask(question: str = Form(...), focus_file: str = Form(None), u=Depends(auth.get_current_user), db: Session = Depends(database.get_db)):
    a, s = rag_engine.get_answer(u.id, question, focus_file)
    db.add(database.ChatMessage(user_id=u.id, question=question, answer=a, sources=json.dumps(s)))
    db.commit(); return {"answer": a, "sources": s}
@app.get("/chat-history")
def hist(u: database.User = Depends(auth.get_current_user), db: Session = Depends(database.get_db)):
    c = db.query(database.ChatMessage).filter(database.ChatMessage.user_id == u.id).all()
    return [{"question": x.question, "answer": x.answer, "sources": json.loads(x.sources)} for x in c]
@app.delete("/chat-history")
def clear(u: database.User = Depends(auth.get_current_user), db: Session = Depends(database.get_db)):
    db.query(database.ChatMessage).filter(database.ChatMessage.user_id == u.id).delete()
    db.commit(); return {"msg": "ok"}
@app.post("/summarize")
def summarize(u=Depends(auth.get_current_user)): return {"summary": rag_engine.generate_document_summary(u.id)}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)""",

    "frontend/login.html": """<!DOCTYPE html><html><head><title>ResearchAI - Access</title>
<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.0.0/css/all.min.css">
<style>
:root { --bg: #0b0f19; --emerald: #10b981; --dark: #111827; --transition: all 0.6s cubic-bezier(0.68, -0.55, 0.265, 1.55); }
body, html { margin:0; width:100%; height:100%; background:var(--bg); overflow:hidden; font-family: sans-serif; }
#bgCanvas { position:fixed; top:0; left:0; width:100%; height:100%; z-index:1; opacity:0.4; pointer-events:none; }
.wrapper { position:relative; z-index:10; width:100%; height:100%; display:flex; justify-content:center; align-items:center; }
.container { position:relative; width:850px; height:500px; background:rgba(17, 24, 39, 0.8); border-radius:24px; box-shadow:0 0 40px rgba(16, 185, 129, 0.1); overflow:hidden; display:flex; border:1px solid rgba(255,255,255,0.05); backdrop-filter: blur(12px); }
.form-container { width:50%; height:100%; display:flex; flex-direction:column; justify-content:center; padding:40px; transition: var(--transition); }
.overlay-container { position:absolute; top:0; left:50%; width:50%; height:100%; background:linear-gradient(135deg, #064e3b, #065f46); z-index:100; transition: var(--transition); display:flex; align-items:center; justify-content:center; text-align:center; color:white; }
.container.active .overlay-container { transform:translateX(-100%); }
.container.active .login-box { transform:translateX(100%); opacity:0; pointer-events:none; }
input { width:100%; padding:14px; margin:10px 0; background:rgba(255,255,255,0.05); border:1px solid #374151; color:white; border-radius:12px; outline:none; }
.btn { background:var(--emerald); color:white; border:none; padding:14px; border-radius:12px; cursor:pointer; font-weight:bold; width:100%; }
.ghost { background:transparent; border:2px solid #fff; margin-top:20px; width:auto; }
</style></head><body>
<canvas id="bgCanvas"></canvas>
<div class="wrapper"><div class="container" id="box">
    <div class="form-container"><form id="rf"><h1>Join Us</h1><input id="rn" placeholder="Name" required><input id="re" placeholder="Email" required><input id="rp" type="password" placeholder="Password" required><button type="submit" class="btn">REGISTER</button></form></div>
    <div class="form-container login-box"><form id="lf"><h1>Welcome Back</h1><input id="le" placeholder="Email" required><input id="lp" type="password" placeholder="Password" required><button type="submit" class="btn">SIGN IN</button></form></div>
    <div class="overlay-container"><div><h2>ResearchAI</h2><button class="btn ghost" onclick="toggle()">SWITCH</button></div></div>
</div></div>
<script>
const canvas=document.getElementById('bgCanvas'), ctx=canvas.getContext('2d'); let pts=[];
function init(){ canvas.width=innerWidth; canvas.height=innerHeight; pts=[]; for(let i=0;i<100;i++) pts.push({x:Math.random()*canvas.width, y:Math.random()*canvas.height, vx:(Math.random()-0.5)*1.5, vy:(Math.random()-0.5)*1.5}); }
function anim(){ ctx.clearRect(0,0,canvas.width,canvas.height); ctx.fillStyle="#10b981"; ctx.strokeStyle="rgba(16, 185, 129, 0.1)"; pts.forEach((p,i)=>{ p.x+=p.vx; p.y+=p.vy; if(p.x<0||p.x>canvas.width) p.vx*=-1; if(p.y<0||p.y>canvas.height) p.vy*=-1; ctx.beginPath(); ctx.arc(p.x,p.y,2,0,Math.PI*2); ctx.fill(); }); requestAnimationFrame(anim); }
init(); anim();
function toggle(){ document.getElementById('box').classList.toggle('active'); }
document.getElementById('rf').onsubmit=async(e)=>{ e.preventDefault(); const fd=new FormData(); fd.append('name',document.getElementById('rn').value); fd.append('email',document.getElementById('re').value); fd.append('password',document.getElementById('rp').value); await fetch('http://localhost:8000/register',{method:'POST',body:fd}); alert("Success!"); toggle(); };
document.getElementById('lf').onsubmit=async(e)=>{ e.preventDefault(); const fd=new FormData(); fd.append('email',document.getElementById('le').value); fd.append('password',document.getElementById('lp').value); const r=await fetch('http://localhost:8000/login',{method:'POST',body:fd}); if(r.ok){ const d=await r.json(); localStorage.setItem('token', d.access_token); window.location.href='app.html'; } else alert("Failed"); };
</script></body></html>""",

    "frontend/app.html": """<!DOCTYPE html><html><head><title>Dashboard</title>
<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.0.0/css/all.min.css">
<script src="https://cdn.jsdelivr.net/npm/marked/marked.min.js"></script>
<style>
:root { --bg: #0b0f19; --emerald: #10b981; --dark: #111827; --card: rgba(31, 41, 55, 0.7); --border: rgba(16, 185, 129, 0.2); }
body { margin:0; display:flex; height:100vh; background:var(--bg); color:#fff; font-family: sans-serif; overflow:hidden; }
.sidebar { width:260px; background:var(--dark); border-right:1px solid var(--border); display:flex; flex-direction:column; }
.main { flex-grow:1; display:flex; flex-direction:column; overflow:hidden; }
.grid { display:grid; grid-template-columns: 360px 1fr; gap:20px; padding:20px; height:100%; overflow:hidden; }
.pane { background:var(--card); border-radius:20px; border:1px solid var(--border); display:flex; flex-direction:column; overflow:hidden; backdrop-filter:blur(12px); }
.pane:hover { border-color:var(--emerald); box-shadow:0 0 30px rgba(16, 185, 129, 0.1); }
#msgs { flex-grow:1; padding:24px; overflow-y:auto; display:flex; flex-direction:column; gap:16px; }
.msg { padding:14px; border-radius:18px; font-size:14px; line-height:1.6; max-width:85%; }
.ai { align-self:flex-start; background:rgba(255,255,255,0.05); }
.user { align-self:flex-end; background:var(--emerald); }
.input-bar { padding:20px; background:rgba(0,0,0,0.3); border-top:1px solid var(--border); display:flex; gap:10px; }
input { flex-grow:1; background:#000; border:1px solid #374151; padding:12px; color:white; border-radius:12px; outline:none; }
.btn { background:var(--emerald); color:white; border:none; padding:10px 20px; border-radius:10px; cursor:pointer; font-weight:bold; }
.source-tag { display:block; margin-top:8px; color:var(--emerald); font-size:11px; cursor:pointer; text-decoration:underline; }
.modal { display:none; position:fixed; z-index:999; left:0; top:0; width:100%; height:100%; background:rgba(0,0,0,0.8); align-items:center; justify-content:center; }
.modal-content { background:var(--dark); padding:30px; border-radius:20px; max-width:600px; width:90%; border:1px solid var(--emerald); }
</style></head><body>
<div class="sidebar"><div style="padding:30px; font-size:22px; font-weight:900; color:var(--emerald)">RESEARCH AI</div><div style="padding:20px; cursor:pointer" onclick="localStorage.clear();location.href='login.html'">Logout</div></div>
<div class="main"><div class="grid">
    <div class="pane" style="padding:20px"><input type="file" id="fi" multiple onchange="up()"><div id="f-l" style="margin-top:20px; overflow-y:auto"></div></div>
    <div class="pane"><div style="padding:15px; display:flex; justify-content:space-between"><button class="btn" id="sumBtn" onclick="sum()">AI Summary</button></div>
    <div id="msgs"></div><div class="input-bar"><input id="qi" placeholder="Ask AI..." onkeypress="if(event.key=='Enter') ask()"><button class="btn" onclick="ask()">SEND</button></div></div>
</div></div>
<div id="sm" class="modal" onclick="if(event.target==this) this.style.display='none'"><div class="modal-content"><h4 id="sm-h" style="color:var(--emerald)">Source</h4><div id="sm-b" style="font-size:14px; line-height:1.6; max-height:300px; overflow-y:auto"></div></div></div>
<script>
const h={'Authorization':'Bearer '+localStorage.getItem('token')}; let focus=null;
async function smoothType(el, txt){ el.innerHTML=""; const words=txt.split(" "); for(let w of words){ el.innerHTML+=w+" "; document.getElementById('msgs').scrollTop=99999; await new Promise(r=>setTimeout(r,40)); } el.innerHTML=marked.parse(txt); }
async function init(){ const r=await fetch('http://localhost:8000/me',{headers:h}); const u=await r.json(); loadDocs(); loadHistory(); }
async function loadHistory(){ const r=await fetch('http://localhost:8000/chat-history',{headers:h}); const hist=await r.json(); document.getElementById('msgs').innerHTML=""; hist.forEach(c=>{ append('user', c.question); append('ai', c.answer, c.sources); }); }
function append(t, txt, s=[]){ const m=document.getElementById('msgs'), d=document.createElement('div'); d.className=`msg ${t}`; let srcHtml=""; if(s.length){ const safeTxt=encodeURIComponent(s[0].text); srcHtml=`<span class="source-tag" onclick="showModal('${s[0].filename}',${s[0].page},'${safeTxt}')">Source: ${s[0].filename} (p${s[0].page})</span>`; } d.innerHTML=(t==='ai'?marked.parse(txt):txt)+srcHtml; m.appendChild(d); m.scrollTop=99999; return d; }
async function ask(){ const i=document.getElementById('qi'), q=i.value; if(!q) return; i.value=""; append('user',q); const ad=append('ai','Thinking...'); const fd=new FormData(); fd.append('question',q); if(focus) fd.append('focus_file',focus); const r=await fetch('http://localhost:8000/ask',{method:'POST',headers:h,body:fd}); const d=await r.json(); await smoothType(ad, d.answer); }
async function loadDocs(){ const r=await fetch('http://localhost:8000/documents',{headers:h}); const docs=await r.json(); const l=document.getElementById('f-l'); l.innerHTML=""; docs.forEach(d=>{ l.innerHTML += `<div style="background:rgba(255,255,255,0.03); padding:10px; border-radius:10px; margin-bottom:10px; display:flex; justify-content:space-between"><span>${d.filename}</span><div style="display:flex;gap:5px"><button class="btn" style="padding:4px 8px;font-size:10px" onclick="focus='${d.filename}';document.getElementById('cur-f').innerText=focus">FOCUS</button><i class="fas fa-trash" style="cursor:pointer" onclick="del(${d.id})"></i></div></div>`; }); }
async function up(){ const fd=new FormData(); const files=document.getElementById('fi').files; for(let i=0; i<files.length; i++) fd.append('files', files[i]); await fetch('http://localhost:8000/upload',{method:'POST',headers:h,body:fd}); loadDocs(); }
async function del(id){ if(confirm("Delete?")){ await fetch(`http://localhost:8000/documents/${id}`,{method:'DELETE',headers:h}); loadDocs(); } }
async function sum(){ const btn=document.getElementById('sumBtn'); btn.innerText='Generating...'; const r=await fetch('http://localhost:8000/summarize',{method:'POST',headers:h}); const d=await r.json(); await smoothType(append('ai','...'), d.summary); btn.innerText='AI Summary'; }
function showModal(f,p,t){ document.getElementById('sm-h').innerText=`${f} - Page ${p}`; document.getElementById('sm-b').innerText=decodeURIComponent(t); document.getElementById('sm').style.display='flex'; }
init();
</script></body></html>"""
}

# 1. Create folders
for folder in structure:
    os.makedirs(os.path.join(base_dir, folder), exist_ok=True)

# 2. Create files
for path, content in files.items():
    with open(os.path.join(base_dir, path), "w", encoding="utf-8") as f:
        f.write(content.strip())

# 3. Create PPT
try:
    from pptx import Presentation
    prs = Presentation()
    slides = [
        ("ResearchAI", "AI-Powered Academic Assistant\\nA Final Year CS Project"),
        ("Problem", "Hallucinations and lack of local knowledge in modern AI."),
        ("Solution", "RAG (Retrieval Augmented Generation) using FAISS and Groq Llama 3."),
        ("Architecture", "FastAPI Backend + HTML5/CSS3/JS Frontend."),
        ("UI Design", "Futuristic Emerald Glassmorphism Theme."),
        ("Key Features", "OCR Support, Multi-format Parsing, Persistent History.")
    ]
    for title, text in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = text
    prs.save(os.path.join(base_dir, "ResearchAI_Presentation.pptx"))
    print("Presentation Generated!")
except ImportError:
    print("Could not generate PPT (python-pptx missing). Code files created successfully.")

print(f"✅ Success! All files created in the '{base_dir}' folder.")
print("👉 Step 1: pip install -r requirements.txt")
print("👉 Step 2: Set your GROQ_API_KEY in .env")
print("👉 Step 3: cd backend && python main.py")